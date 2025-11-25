# utils/parser_utils.py
from pathlib import Path
from typing import Tuple
import pdfplumber, PyPDF2, docx, pandas as pd
from PIL import Image
import pytesseract
from io import BytesIO

MAX_CHARS_RETURN = 5000

def parse_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return text

def parse_pdf(path: str, max_pages: int = 0) -> str:
    text = []
    try:
        with pdfplumber.open(path) as pdf:
            pages = pdf.pages
            for i, p in enumerate(pages):
                if max_pages and i >= max_pages:
                    break
                t = p.extract_text()
                if t:
                    text.append(t)
    except Exception:
        # fallback to PyPDF2 simple extraction
        try:
            with open(path, "rb") as fh:
                reader = PyPDF2.PdfReader(fh)
                for i, page in enumerate(reader.pages):
                    if max_pages and i >= max_pages:
                        break
                    try:
                        text.append(page.extract_text() or "")
                    except:
                        continue
        except Exception as e:
            return f"[PDF parse error] {e}"
    out = "\n".join(text)
    return out

def parse_docx(path: str) -> str:
    doc = docx.Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)

def parse_xlsx(path: str) -> str:
    df = pd.read_excel(path, sheet_name=0)
    return df.astype(str).to_string()

def parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    lines.append(txt)
    return "\n".join(lines)

def parse_image_ocr(path: str) -> str:
    try:
        im = Image.open(path)
        text = pytesseract.image_to_string(im, lang='chi_sim+eng')
        return text
    except Exception as e:
        return f"[OCR Error] {e}"

def parse_file(path: str, max_pages: int = 0) -> str:
    p = Path(path)
    ext = p.suffix.lower()
    if ext in [".txt", ".md", ".py", ".json", ".csv"]:
        return parse_txt(path)
    if ext == ".pdf":
        return parse_pdf(path, max_pages=max_pages)
    if ext in [".docx"]:
        return parse_docx(path)
    if ext in [".xlsx", ".xls", ".csv"]:
        return parse_xlsx(path)
    if ext in [".pptx", ".ppt"]:
        return parse_pptx(path)
    if ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
        return parse_image_ocr(path)
    return f"[Unsupported file type: {ext}]"

def get_preview(text: str, n_chars: int = 800) -> str:
    if not text:
        return ""
    return text[:n_chars] + ("…(truncated)" if len(text) > n_chars else "")
